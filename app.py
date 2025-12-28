import streamlit as st
from dotenv import load_dotenv, find_dotenv
from PyPDF2 import PdfReader
from langchain_text_splitters import CharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.messages import HumanMessage, AIMessage
from langchain_core.documents import Document
from htmlTemplates import css, bot_template, user_template
import os
from pathlib import Path
from google import genai
from sentence_transformers import CrossEncoder
import time
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_JUSTIFY
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import re

def get_pdf_text_with_metadata(pdf_docs):
    """Extract text with metadata (page numbers, source)"""
    documents = []
    for pdf_idx, pdf in enumerate(pdf_docs):
        pdf_reader = PdfReader(pdf)
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text.strip():
                doc = Document(
                    page_content=text,
                    metadata={
                        'source': pdf.name,
                        'page': page_num + 1,
                        'pdf_index': pdf_idx
                    }
                )
                documents.append(doc)
    return documents

def get_pdf_info(pdf_file):
    """Get basic PDF information"""
    try:
        pdf_file.seek(0)
        pdf_reader = PdfReader(pdf_file)
        return {
            'pages': len(pdf_reader.pages),
            'name': pdf_file.name,
            'size': f"{pdf_file.size / 1024:.1f} KB"
        }
    except Exception as e:
        return None

def get_text_chunks_with_metadata(documents):
    """Split documents into chunks while preserving metadata"""
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = []
    for doc in documents:
        splits = text_splitter.split_text(doc.page_content)
        for split in splits:
            chunk = Document(
                page_content=split,
                metadata=doc.metadata
            )
            chunks.append(chunk)
    return chunks

def get_vectorstore(chunks):
    """Create vector store with metadata"""
    from langchain_huggingface import HuggingFaceEmbeddings
    hf_token = os.getenv("HUGGINGFACEHUB_API_TOKEN")
    if hf_token:
        os.environ["HUGGINGFACEHUB_API_TOKEN"] = hf_token
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vectorstore = FAISS.from_documents(chunks, embeddings)
    return vectorstore

def format_docs(docs):
    """Format documents for display"""
    return "\n\n".join([f"[Source: {doc.metadata.get('source', 'Unknown')}, Page {doc.metadata.get('page', '?')}]\n{doc.page_content}" for doc in docs])

def enhance_query(original_query):
    """Improve user's question for better retrieval"""
    try:
        client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
        prompt = f"""You are an expert at improving search queries to find the most relevant information in documents.

Task: Transform the user's question into an optimized search query that will retrieve the best matching content.

Original question: "{original_query}"

Instructions:
1. Identify the core intent and key concepts
2. Expand abbreviations and add relevant synonyms
3. Make the query more specific and detailed
4. Keep it as a clear, natural language question
5. Focus on searchable keywords

Enhanced question (only output the improved question, nothing else):"""
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            enhanced = response.text.strip()
            keywords = [word for word in original_query.split() if len(word) > 4][:5]
            return enhanced, keywords
    except Exception as e:
        pass
    return original_query, []

def generate_multi_queries(question):
    """Generate multiple query variations"""
    try:
        client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
        prompt = f"""You are an expert at generating diverse search queries to maximize information retrieval.

Original question: "{question}"

Task: Generate 2 alternative versions of this question that:
1. Use different wording but maintain the same intent
2. Approach the topic from different angles
3. Include synonyms and related terms
4. Are specific and clear

Output only the 2 alternative questions, one per line, without numbering or labels:"""
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            queries = [question]
            for line in response.text.split('\n'):
                line = line.strip()
                if line and len(line) > 10:
                    cleaned = line
                    for prefix in ['Alternative 1:', 'Alternative 2:', '1.', '2.', '1)', '2)', 'Alternative:', '-']:
                        cleaned = cleaned.replace(prefix, '').strip()
                    if cleaned and cleaned not in queries and len(cleaned) > 10:
                        queries.append(cleaned)
            return queries[:3]
    except Exception as e:
        pass
    return [question]

def rerank_documents(query, documents, top_k=3):
    """Re-rank documents using cross-encoder"""
    if not documents:
        return documents
    try:
        model = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
        pairs = [[query, doc.page_content] for doc in documents]
        scores = model.predict(pairs)
        ranked = sorted(zip(documents, scores), key=lambda x: x[1], reverse=True)
        return [doc for doc, score in ranked[:top_k]]
    except Exception as e:
        return documents[:top_k]

def query_gemini(context, question, max_retries=2):
    """Query Gemini with enhanced prompt for better answers"""
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
    models_to_try = ['gemini-2.5-flash', 'gemini-2.0-flash', 'gemini-flash-latest']
    if 'working_gemini_model' in st.session_state and st.session_state.working_gemini_model:
        models_to_try.insert(0, st.session_state.working_gemini_model)
    for model_name in models_to_try:
        for attempt in range(max_retries):
            try:
                max_context = 4000
                if len(context) > max_context:
                    context = context[:max_context] + "\n\n[Context truncated...]"
                prompt = f"""You are an intelligent assistant specialized in analyzing documents and providing accurate, helpful answers.

CONTEXT FROM DOCUMENTS:
{context}

USER QUESTION:
{question}

INSTRUCTIONS:
1. Read the context carefully and identify all relevant information
2. Answer the question directly and accurately based ONLY on the provided context
3. If the context contains the answer:
   - Provide a clear, comprehensive response
   - Include specific details, numbers, dates, or facts from the context
   - Structure your answer logically (use paragraphs if needed)
   - Be concise but complete
4. If the context does NOT contain enough information:
   - Clearly state: "Based on the provided documents, I cannot find specific information about [topic]."
   - If there's partial information, share what you found and explain what's missing
5. If asked to compare, list, or analyze:
   - Organize information clearly
   - Use bullet points only if it improves clarity
   - Provide complete explanations
6. Cite specific details from the context when possible (e.g., "According to page X...")
7. Be professional, accurate, and helpful

ANSWER:"""
                response = client.models.generate_content(model=model_name, contents=prompt)
                if response and hasattr(response, 'text') and response.text:
                    if 'working_gemini_model' not in st.session_state:
                        st.session_state.working_gemini_model = model_name
                    return response.text
            except Exception as e:
                error_msg = str(e)
                if "429" in error_msg or "quota" in error_msg.lower():
                    if attempt < max_retries - 1:
                        time.sleep(2)
                        continue
                if "404" in error_msg or "not found" in error_msg.lower():
                    break
                if attempt < max_retries - 1:
                    time.sleep(1)
                    continue
    return "I'm having trouble generating a response. Please try again."

def extract_full_document_text(pdf_docs):
    """Extract complete text from all PDFs"""
    full_text = ""
    doc_map = {}
    for pdf in pdf_docs:
        pdf.seek(0)
        pdf_reader = PdfReader(pdf)
        doc_text = ""
        for page in pdf_reader.pages:
            doc_text += page.extract_text() + "\n"
        doc_map[pdf.name] = doc_text
        full_text += f"\n\n=== {pdf.name} ===\n\n{doc_text}"
    return full_text, doc_map

def transform_document(original_text, user_instruction):
    """Transform document based on user instructions using AI"""
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
    prompt = f"""You are an expert document editor and transformer.

ORIGINAL DOCUMENT:
{original_text[:8000]}

USER INSTRUCTION:
{user_instruction}

TASK:
Transform the document according to the user's instruction. This could include:
- Reformatting (change structure, layout, organization)
- Rewriting (improve language, clarity, tone)
- Summarizing (condense content)
- Expanding (add more details)
- Updating (change specific information)
- Translating (convert to another language)
- Extracting (pull out specific sections)
- Combining (merge information)

OUTPUT REQUIREMENTS:
1. Provide the COMPLETE transformed document
2. Maintain professional formatting
3. Keep factual accuracy
4. Preserve important information unless instructed otherwise
5. Use clear section headers if appropriate

TRANSFORMED DOCUMENT:"""
    try:
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            return response.text
    except Exception as e:
        return f"Error transforming document: {str(e)}"
    return "Unable to transform document"

def generate_pdf(content, filename="transformed_document.pdf"):
    """Generate a new PDF from text content"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='Justify', parent=styles['BodyText'], alignment=TA_JUSTIFY, fontSize=11, leading=14))
    story = []
    paragraphs = content.split('\n\n')
    for para_text in paragraphs:
        para_text = para_text.strip()
        if not para_text:
            continue
        if para_text.startswith('#'):
            para_text = para_text.replace('#', '').strip()
            para = Paragraph(para_text, styles['Heading1'])
        elif para_text.isupper() or para_text.endswith(':'):
            para = Paragraph(para_text, styles['Heading2'])
        else:
            para = Paragraph(para_text, styles['Justify'])
        story.append(para)
        story.append(Spacer(1, 0.2 * inch))
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def compare_documents(original, transformed):
    """Compare original and transformed documents"""
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
    prompt = f"""Compare these two document versions and highlight the changes:

ORIGINAL (first 2000 chars):
{original[:2000]}

TRANSFORMED (first 2000 chars):
{transformed[:2000]}

Provide:
1. **Key Changes**: What was modified?
2. **Additions**: What was added?
3. **Removals**: What was removed?
4. **Structure Changes**: How is the organization different?
5. **Summary**: Overall transformation description

Keep it concise (5-7 bullet points total)."""
    try:
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            return response.text
    except:
        pass
    return "Comparison not available"

def generate_presentation_content(document_text, num_slides=10):
    """Generate presentation structure from document"""
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
    prompt = f"""You are a presentation expert. Create a professional PowerPoint presentation structure from this document.

DOCUMENT TEXT:
{document_text[:6000]}

REQUIREMENTS:
1. Create {num_slides} slides maximum (including title slide)
2. Each slide should have:
   - A clear, concise title
   - 3-5 bullet points OR key content
   - Brief, impactful text (not paragraphs)

OUTPUT FORMAT (JSON):
{{
  "title": "Main Presentation Title",
  "subtitle": "Brief subtitle or topic",
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
      "notes": "Optional speaker notes"
    }}
  ]
}}

Generate the presentation structure in valid JSON format:"""
    
    try:
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
    except Exception as e:
        st.error(f"Error generating presentation: {e}")
    return None

def create_powerpoint(presentation_data):
    """Create PowerPoint from presentation data"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = presentation_data.get('title', 'Document Presentation')
    subtitle.text = presentation_data.get('subtitle', 'Generated from PDF')
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content Slides
    for slide_data in presentation_data.get('slides', []):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = slide_data.get('title', 'Untitled')
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for bullet_point in slide_data.get('content', []):
            p = text_frame.add_paragraph()
            p.text = bullet_point
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(6)
        
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def generate_qa_pairs(document_text, num_pairs=20):
    """Generate Q&A pairs from document"""
    client = genai.Client(api_key=os.getenv("GOOGLE_API_KEY"))
    prompt = f"""You are an expert educator. Generate comprehensive question-answer pairs from this document.

DOCUMENT TEXT:
{document_text[:6000]}

REQUIREMENTS:
1. Generate {num_pairs} question-answer pairs
2. Questions should:
   - Cover key concepts and details
   - Range from basic to advanced
   - Include different types (factual, conceptual, analytical)
   - Be clear and specific
3. Answers should:
   - Be accurate and complete
   - Reference the document content
   - Be 2-4 sentences long

OUTPUT FORMAT (JSON):
{{
  "qa_pairs": [
    {{
      "question": "What is...?",
      "answer": "The answer is...",
      "difficulty": "easy|medium|hard",
      "category": "factual|conceptual|analytical"
    }}
  ]
}}

Generate the Q&A pairs in valid JSON format:"""
    
    try:
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        if response and response.text:
            json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
    except Exception as e:
        st.error(f"Error generating Q&A: {e}")
    return None

def format_qa_as_text(qa_data):
    """Format Q&A pairs as readable text"""
    text = "QUESTION & ANSWER PAIRS\n"
    text += "=" * 80 + "\n\n"
    
    for idx, qa in enumerate(qa_data.get('qa_pairs', []), 1):
        text += f"Q{idx}. {qa.get('question', 'No question')}\n"
        text += f"Difficulty: {qa.get('difficulty', 'N/A').upper()} | "
        text += f"Type: {qa.get('category', 'N/A').upper()}\n\n"
        text += f"A{idx}. {qa.get('answer', 'No answer')}\n\n"
        text += "-" * 80 + "\n\n"
    
    return text

def format_qa_as_pdf(qa_data):
    """Format Q&A pairs as PDF"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    styles = getSampleStyleSheet()
    
    styles.add(ParagraphStyle(name='Question', parent=styles['Heading2'], fontSize=12, textColor='blue', spaceAfter=6))
    styles.add(ParagraphStyle(name='Answer', parent=styles['BodyText'], fontSize=11, spaceAfter=12))
    styles.add(ParagraphStyle(name='Meta', parent=styles['BodyText'], fontSize=9, textColor='gray', spaceAfter=12))
    
    story = []
    story.append(Paragraph("QUESTION & ANSWER PAIRS", styles['Title']))
    story.append(Spacer(1, 0.3 * inch))
    
    for idx, qa in enumerate(qa_data.get('qa_pairs', []), 1):
        story.append(Paragraph(f"Q{idx}. {qa.get('question', 'No question')}", styles['Question']))
        meta = f"Difficulty: {qa.get('difficulty', 'N/A').upper()} | Type: {qa.get('category', 'N/A').upper()}"
        story.append(Paragraph(meta, styles['Meta']))
        story.append(Paragraph(f"A{idx}. {qa.get('answer', 'No answer')}", styles['Answer']))
        story.append(Spacer(1, 0.2 * inch))
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def process_query_advanced(user_question, vectorstore, use_enhancement=True, use_multi_query=True, use_reranking=True):
    """Process query with advanced techniques"""
    results = {'original_query': user_question, 'enhanced_query': user_question, 'keywords': [], 'all_docs': [], 'final_docs': []}
    if use_enhancement:
        try:
            enhanced, keywords = enhance_query(user_question)
            results['enhanced_query'] = enhanced
            results['keywords'] = keywords
        except:
            pass
    queries_to_search = [user_question]
    if use_multi_query:
        try:
            multi_queries = generate_multi_queries(user_question)
            queries_to_search = multi_queries
        except:
            pass
    all_docs = []
    seen_content = set()
    for query in queries_to_search:
        try:
            docs = vectorstore.similarity_search(query, k=5)
            for doc in docs:
                if doc.page_content not in seen_content:
                    all_docs.append(doc)
                    seen_content.add(doc.page_content)
        except:
            continue
    results['all_docs'] = all_docs
    if use_reranking and all_docs:
        try:
            final_docs = rerank_documents(user_question, all_docs, top_k=3)
            results['final_docs'] = final_docs
        except:
            results['final_docs'] = all_docs[:3]
    else:
        results['final_docs'] = all_docs[:3]
    return results

def display_chat_messages():
    """Display all chat messages"""
    if st.session_state.chat_history:
        for message in st.session_state.chat_history:
            if isinstance(message, HumanMessage):
                st.markdown(user_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)
            else:
                st.markdown(bot_template.replace("{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    env_path = find_dotenv()
    if env_path:
        load_dotenv(env_path, override=True)
    else:
        env_file = Path(__file__).parent / ".env"
        if env_file.exists():
            load_dotenv(env_file, override=True)
    st.set_page_config(page_title="Advanced RAG Chat", page_icon="💬", layout="wide")
    st.markdown(css, unsafe_allow_html=True)
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        st.error("⚠️ GOOGLE_API_KEY not found in .env file!")
        st.stop()
    
    # Initialize all session state variables
    session_vars = {
        "vectorstore": None, "chat_history": [], "working_gemini_model": None,
        "pdf_docs": None, "use_enhancement": True, "use_multi_query": True,
        "use_reranking": True, "last_sources": None, "transformed_text": None,
        "transformed_pdf": None, "transformation_comparison": None,
        "presentation_pptx": None, "qa_data": None, "qa_text": None, "qa_pdf": None
    }
    for key, default in session_vars.items():
        if key not in st.session_state:
            st.session_state[key] = default
    
    with st.sidebar:
        st.title("⚙️ Settings")
        with st.expander("🎛️ Advanced Features", expanded=False):
            st.session_state.use_enhancement = st.checkbox("🔍 Query Enhancement", value=st.session_state.use_enhancement, help="AI improves your question for better search")
            st.session_state.use_multi_query = st.checkbox("🔄 Multi-Query", value=st.session_state.use_multi_query, help="Searches with multiple question variations")
            st.session_state.use_reranking = st.checkbox("⚡ Re-ranking", value=st.session_state.use_reranking, help="Re-orders results by relevance")
        
        st.divider()
        st.subheader("📄 Upload Documents")
        pdf_docs = st.file_uploader("Choose PDFs", accept_multiple_files=True, type=['pdf'], label_visibility="collapsed")
        
        if pdf_docs:
            with st.expander("📑 Document Info", expanded=False):
                total_pages = 0
                for pdf in pdf_docs:
                    info = get_pdf_info(pdf)
                    if info:
                        total_pages += info['pages']
                        st.text(f"📄 {info['name'][:30]}... ({info['pages']} pages)")
                st.info(f"Total: {len(pdf_docs)} docs, {total_pages} pages")
        
        if st.button("🚀 Process", type="primary", use_container_width=True):
            if not pdf_docs:
                st.error("Upload PDFs first!")
            else:
                with st.spinner("Processing..."):
                    try:
                        documents = get_pdf_text_with_metadata(pdf_docs)
                        chunks = get_text_chunks_with_metadata(documents)
                        vectorstore = get_vectorstore(chunks)
                        st.session_state.vectorstore = vectorstore
                        st.session_state.pdf_docs = pdf_docs
                        st.session_state.chat_history = []
                        st.success("✅ Ready to chat!")
                        st.balloons()
                    except Exception as e:
                        st.error(f"Error: {str(e)}")
        
        # Document Transformation Section
        if st.session_state.vectorstore is not None:
            st.divider()
            st.subheader("✨ Document Tools")
            
            tool_choice = st.radio(
                "Select Tool:",
                ["📝 Transform Document", "🎨 Create Presentation", "❓ Generate Q&A"],
                label_visibility="collapsed"
            )
            
            # Transform Document
            if tool_choice == "📝 Transform Document":
                transformation_mode = st.selectbox("Transformation Type", 
                    ["💬 Select Option", "📝 Reformat Document", "✍️ Rewrite & Improve", 
                     "📊 Summarize", "🔍 Extract Info", "🌍 Translate", "🎯 Custom"])
                
                if transformation_mode == "🎯 Custom":
                    transform_instruction = st.text_area("Enter instruction:", 
                        placeholder="e.g., 'Convert to cover letter'", height=80)
                else:
                    instructions_map = {
                        "💬 Select Option": None,
                        "📝 Reformat Document": "Reformat with better structure and clear headings.",
                        "✍️ Rewrite & Improve": "Improve clarity, grammar, and professionalism.",
                        "📊 Summarize": "Create concise summary with key points.",
                        "🔍 Extract Info": "Extract important facts and data.",
                        "🌍 Translate": "Translate to English (specify language if needed)."
                    }
                    transform_instruction = instructions_map.get(transformation_mode, "")
                
                additional_notes = st.text_input("Additional notes:", placeholder="Optional requirements")
                if additional_notes and transform_instruction:
                    transform_instruction += f"\n\n{additional_notes}"
                
                if st.button("🚀 Transform", use_container_width=True, disabled=(transformation_mode == "💬 Select Option")):
                    if not transform_instruction:
                        st.error("Provide instructions!")
                    else:
                        with st.spinner("Transforming..."):
                            try:
                                full_text, _ = extract_full_document_text(st.session_state.pdf_docs)
                                transformed_text = transform_document(full_text, transform_instruction)
                                comparison = compare_documents(full_text[:3000], transformed_text[:3000])
                                pdf_bytes = generate_pdf(transformed_text)
                                st.session_state.transformed_text = transformed_text
                                st.session_state.transformed_pdf = pdf_bytes
                                st.session_state.transformation_comparison = comparison
                                st.success("✅ Transformed!")
                                st.balloons()
                            except Exception as e:
                                st.error(f"Error: {str(e)}")
                
                if st.session_state.get('transformed_text'):
                    st.download_button("📄 Download PDF", data=st.session_state.transformed_pdf, 
                        file_name="transformed.pdf", mime="application/pdf", use_container_width=True)
                    st.download_button("📝 Download Text", data=st.session_state.transformed_text, 
                        file_name="transformed.txt", mime="text/plain", use_container_width=True)
                    with st.expander("📊 Changes"):
                        st.markdown(st.session_state.transformation_comparison)
            
            # Create Presentation
            elif tool_choice == "🎨 Create Presentation":
                num_slides = st.slider("Number of slides:", 5, 20, 10)
                presentation_style = st.selectbox("Style:", 
                    ["Professional", "Academic", "Creative", "Minimalist"])
                
                if st.button("🎨 Generate Presentation", use_container_width=True, type="primary"):
                    with st.spinner("Creating presentation..."):
                        try:
                            full_text, _ = extract_full_document_text(st.session_state.pdf_docs)
                            
                            with st.status("🎯 Analyzing content...", expanded=True) as status:
                                pres_data = generate_presentation_content(full_text, num_slides)
                                st.write("✅ Content structured!")
                                status.update(label="✅ Analysis complete!", state="complete")
                            
                            if pres_data:
                                with st.status("🎨 Creating slides...", expanded=True) as status:
                                    pptx_bytes = create_powerpoint(pres_data)
                                    st.write("✅ Presentation created!")
                                    status.update(label="✅ Slides ready!", state="complete")
                                
                                st.session_state.presentation_pptx = pptx_bytes
                                st.success(f"🎉 Created {len(pres_data.get('slides', []))+1} slides!")
                                st.balloons()
                            else:
                                st.error("Failed to generate presentation structure")
                        except Exception as e:
                            st.error(f"Error: {str(e)}")
                
                if st.session_state.get('presentation_pptx'):
                    st.download_button("📊 Download PowerPoint", 
                        data=st.session_state.presentation_pptx,
                        file_name="presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True, type="primary")
            
            # Generate Q&A
            elif tool_choice == "❓ Generate Q&A":
                num_questions = st.slider("Number of Q&A pairs:", 5, 50, 20)
                difficulty_filter = st.multiselect("Include difficulties:", 
                    ["Easy", "Medium", "Hard"], default=["Easy", "Medium", "Hard"])
                
                if st.button("❓ Generate Q&A", use_container_width=True, type="primary"):
                    with st.spinner("Generating questions..."):
                        try:
                            full_text, _ = extract_full_document_text(st.session_state.pdf_docs)
                            
                            with st.status("🤔 Generating questions...", expanded=True) as status:
                                qa_data = generate_qa_pairs(full_text, num_questions)
                                st.write(f"✅ Generated {len(qa_data.get('qa_pairs', []))} pairs!")
                                status.update(label="✅ Q&A generated!", state="complete")
                            
                            if qa_data:
                                with st.status("📝 Formatting outputs...", expanded=True) as status:
                                    qa_text = format_qa_as_text(qa_data)
                                    qa_pdf = format_qa_as_pdf(qa_data)
                                    st.write("✅ Formatted!")
                                    status.update(label="✅ Ready to download!", state="complete")
                                
                                st.session_state.qa_data = qa_data
                                st.session_state.qa_text = qa_text
                                st.session_state.qa_pdf = qa_pdf
                                st.success(f"🎉 Generated {len(qa_data['qa_pairs'])} Q&A pairs!")
                                st.balloons()
                            else:
                                st.error("Failed to generate Q&A pairs")
                        except Exception as e:
                            st.error(f"Error: {str(e)}")
                
                if st.session_state.get('qa_data'):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.download_button("📄 Download PDF", 
                            data=st.session_state.qa_pdf,
                            file_name="qa_pairs.pdf",
                            mime="application/pdf",
                            use_container_width=True)
                    with col2:
                        st.download_button("📝 Download TXT", 
                            data=st.session_state.qa_text,
                            file_name="qa_pairs.txt",
                            mime="text/plain",
                            use_container_width=True)
                    
                    with st.expander("👀 Preview Q&A", expanded=False):
                        qa_preview = st.session_state.qa_data.get('qa_pairs', [])[:3]
                        for idx, qa in enumerate(qa_preview, 1):
                            st.markdown(f"**Q{idx}:** {qa.get('question')}")
                            st.caption(f"*{qa.get('difficulty', 'N/A')} | {qa.get('category', 'N/A')}*")
                            st.markdown(f"**A{idx}:** {qa.get('answer')}")
                            st.divider()
        
        if st.session_state.chat_history:
            st.divider()
            if st.button("🗑️ Clear Chat", use_container_width=True):
                st.session_state.chat_history = []
                st.session_state.last_sources = None
                st.rerun()
        
        if st.session_state.last_sources:
            st.divider()
            with st.expander("📚 Last Sources", expanded=False):
                for i, doc in enumerate(st.session_state.last_sources, 1):
                    st.caption(f"**{i}.** {doc.metadata.get('source', 'Unknown')} - Page {doc.metadata.get('page', '?')}")
    
    # Main chat area
    st.title("💬 Chat with Your PDFs")
    display_chat_messages()
    if not st.session_state.chat_history:
        st.info("👋 Upload PDFs in the sidebar and start chatting!")
    
    user_question = st.chat_input("Ask about your documents...")
    if user_question:
        if st.session_state.vectorstore is None:
            st.warning("⚠️ Please upload and process PDFs first!")
        else:
            st.session_state.chat_history.append(HumanMessage(content=user_question))
            st.markdown(user_template.replace("{{MSG}}", user_question), unsafe_allow_html=True)
            try:
                with st.spinner("🤔 Thinking..."):
                    results = process_query_advanced(user_question, st.session_state.vectorstore, 
                        st.session_state.use_enhancement, st.session_state.use_multi_query, st.session_state.use_reranking)
                    if not results['final_docs']:
                        response = "I couldn't find relevant information in the documents."
                    else:
                        context = format_docs(results['final_docs'])
                        response = query_gemini(context, user_question)
                        st.session_state.last_sources = results['final_docs']
                    st.session_state.chat_history.append(AIMessage(content=response))
                    st.markdown(bot_template.replace("{{MSG}}", response), unsafe_allow_html=True)
                    st.rerun()
            except Exception as e:
                error_msg = f"Error: {str(e)}"
                st.session_state.chat_history.append(AIMessage(content=error_msg))
                st.error(error_msg)
                st.rerun()

if __name__ == '__main__':
    main()
